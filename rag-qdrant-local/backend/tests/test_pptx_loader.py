"""PPTX loader — slide-per-segment, page numbered like PDF."""

from __future__ import annotations

from pathlib import Path

import pytest

from app.document_loader import DocumentLoadError, load_document

pytest.importorskip("pptx")


def _build_pptx(path: Path) -> None:
    """Synthesize a small PPTX with two slides, one with a table + speaker notes."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    # Slide 1 — title + body
    layout = prs.slide_layouts[1]
    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Firewall Konzept"
    slide1.placeholders[1].text = (
        "Wir nutzen pfSense zum Schutz des Verwaltungsnetzes."
    )

    # Slide 2 — title + table + speaker notes
    layout2 = prs.slide_layouts[5]  # title only
    slide2 = prs.slides.add_slide(layout2)
    slide2.shapes.title.text = "Regelwerk"

    rows, cols = 2, 2
    left = top = Inches(1)
    width = Inches(4)
    height = Inches(1.5)
    table = slide2.shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "Quelle"
    table.cell(0, 1).text = "Ziel"
    table.cell(1, 0).text = "LAN"
    table.cell(1, 1).text = "Internet"

    slide2.notes_slide.notes_text_frame.text = "Review im Q2 2026."

    prs.save(str(path))


def test_pptx_loader_emits_one_segment_per_slide(tmp_path: Path):
    sample = tmp_path / "deck.pptx"
    _build_pptx(sample)

    doc = load_document(sample)
    assert doc.file_extension == ".pptx"
    assert doc.document_type == "pptx"
    assert len(doc.segments) == 2
    assert [s.page for s in doc.segments] == [1, 2]


def test_pptx_loader_captures_table_cells_and_speaker_notes(tmp_path: Path):
    sample = tmp_path / "deck.pptx"
    _build_pptx(sample)

    doc = load_document(sample)
    second = doc.segments[1]
    # Header row appears
    assert "Quelle | Ziel" in second.text
    # Data row appears
    assert "LAN | Internet" in second.text
    # Speaker notes flagged with [Notizen]
    assert "[Notizen]" in second.text
    assert "Q2 2026" in second.text


def test_pptx_loader_keeps_slide_one_body(tmp_path: Path):
    sample = tmp_path / "deck.pptx"
    _build_pptx(sample)

    doc = load_document(sample)
    first = doc.segments[0]
    assert "Firewall Konzept" in first.text
    assert "pfSense" in first.text


def test_pptx_loader_raises_for_empty_deck(tmp_path: Path):
    """A presentation with no text at all must surface as DocumentLoadError."""
    from pptx import Presentation
    p = tmp_path / "empty.pptx"
    Presentation().save(str(p))

    with pytest.raises(DocumentLoadError):
        load_document(p)
